from __future__ import annotations

import json
import subprocess
import sys
import tempfile
import unittest
from pathlib import Path
from unittest import mock


SCRIPTS = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(SCRIPTS))

import source_pipeline as pipeline


class MockConverter:
    name = "mock-markitdown"
    version = "1.2.3"

    def __init__(self) -> None:
        self.calls: list[Path] = []

    def convert(
        self,
        source: Path,
        destination: Path,
        config: dict[str, object],
    ) -> None:
        self.calls.append(source)
        if source.name == "broken.docx":
            raise ValueError("malformed document")
        destination.write_text(
            f"# Derived\n\nConverted: {source.read_text(encoding='utf-8')}",
            encoding="utf-8",
        )


class SourceFixture(unittest.TestCase):
    def setUp(self) -> None:
        self.temp = tempfile.TemporaryDirectory(prefix="source-pipeline-")
        self.root = Path(self.temp.name)
        subprocess.run(["git", "init", "-q"], cwd=self.root, check=True)
        subprocess.run(
            ["git", "config", "user.email", "fixture@example.test"],
            cwd=self.root,
            check=True,
        )
        subprocess.run(
            ["git", "config", "user.name", "Fixture"],
            cwd=self.root,
            check=True,
        )
        self.wiki = self.root / "docs" / "wiki"
        (self.wiki / "_system" / "generated" / "sources").mkdir(
            parents=True
        )
        (self.wiki / "_system" / "templates").mkdir(parents=True)
        (self.wiki / "_system" / "templates" / "raw-group-pointer.md").write_text(
            pipeline.DEFAULT_POINTER_TEMPLATE,
            encoding="utf-8",
        )

    def tearDown(self) -> None:
        self.temp.cleanup()

    def write(self, relative: str, content: str | bytes = "content") -> Path:
        path = self.root / relative
        path.parent.mkdir(parents=True, exist_ok=True)
        if isinstance(content, bytes):
            path.write_bytes(content)
        else:
            path.write_text(content, encoding="utf-8")
        return path

    def track(self, *relative: str) -> None:
        subprocess.run(["git", "add", "--", *relative], cwd=self.root, check=True)

    def scan(
        self,
        *,
        host: dict[str, object] | None = None,
        previous: dict[str, object] | None = None,
        converter: MockConverter | None = None,
        generated_at: str = "2026-09-05T00:00:00Z",
    ) -> dict[str, object]:
        return pipeline.SourceInventory(
            self.root,
            host=host,
            cache_root=self.wiki / "_system" / "generated" / "source-cache",
            converter=converter,
        ).scan(previous=previous, generated_at=generated_at)


class ConfigurationAndSecurityTests(SourceFixture):
    def test_default_scan_is_git_tracked_and_respects_git_ignores(self) -> None:
        self.write(".gitignore", "ignored/\n")
        self.write("docs/tracked.md", "tracked")
        self.write("docs/untracked.md", "untracked")
        self.write("ignored/dependency.md", "ignored")
        self.track(".gitignore", "docs/tracked.md")

        manifest = self.scan()

        self.assertEqual(
            [entry["path"] for entry in manifest["entries"]],
            [".gitignore", "docs/tracked.md"],
        )

    def test_host_includes_add_untracked_files_and_excludes_win(self) -> None:
        self.write("odd/handbook.note", "include me")
        self.write("docs/keep.md", "keep")
        self.write("docs/drop.md", "drop")
        self.track("docs/keep.md", "docs/drop.md")
        host = {
            "sources": {
                "include": ["odd/**/*.note"],
                "exclude": ["docs/drop.md"],
            }
        }

        config = pipeline.load_source_config(host)
        manifest = self.scan(host=host)

        self.assertEqual(config.include, ("odd/**/*.note",))
        self.assertEqual(
            [entry["path"] for entry in manifest["entries"]],
            ["docs/keep.md", "odd/handbook.note"],
        )

    def test_likely_secrets_are_denied_even_when_explicitly_included(self) -> None:
        self.write(".env", "TOKEN=secret")
        self.write("keys/server.pem", "private")
        self.write("config/credentials.json", "{}")
        self.write("config/settings.yaml", "safe: true")
        self.track(
            ".env",
            "keys/server.pem",
            "config/credentials.json",
            "config/settings.yaml",
        )

        manifest = self.scan(
            host={"sources": {"include": [".env", "keys/server.pem"]}}
        )

        self.assertEqual(
            [entry["path"] for entry in manifest["entries"]],
            ["config/settings.yaml"],
        )
        denied = {item["path"] for item in manifest["denied"]}
        self.assertEqual(
            denied,
            {".env", "config/credentials.json", "keys/server.pem"},
        )

    def test_source_and_cache_paths_cannot_escape_repository(self) -> None:
        outside = self.root.parent / "outside-source.md"
        outside.write_text("outside", encoding="utf-8")
        self.addCleanup(outside.unlink, missing_ok=True)

        with self.assertRaisesRegex(pipeline.SourcePathError, "repository"):
            pipeline.SourceInventory(
                self.root,
                host={"sources": {"include": ["../outside-source.md"]}},
            ).scan()
        with self.assertRaisesRegex(pipeline.SourcePathError, "cache"):
            pipeline.SourceInventory(
                self.root,
                cache_root=self.root.parent / "escaped-cache",
            )


class ClassificationAndStructureTests(SourceFixture):
    def test_classifies_supported_source_categories(self) -> None:
        fixtures = {
            "src/lib.rs": "code",
            "tests/test_api.py": "tests",
            "docs/guide.md": "docs",
            "docs/adr/0001-choice.md": "adr",
            "config/app.toml": "config",
            "data/rows.csv": "data",
            "assets/logo.png": "media-binary",
        }
        for path in fixtures:
            self.write(path, b"\x89PNG\r\n" if path.endswith(".png") else path)
        self.track(*fixtures)

        manifest = self.scan()
        actual = {
            entry["path"]: entry["classification"]
            for entry in manifest["entries"]
        }

        self.assertEqual(actual, fixtures)
        code = next(e for e in manifest["entries"] if e["path"] == "src/lib.rs")
        self.assertEqual(code["conversion"]["state"], "graphify-delegate")

    def test_detects_context_adr_and_common_documentation_sites(self) -> None:
        files = {
            "CONTEXT.md": "# Context",
            "docs/adr/0001-use-sql.md": "# ADR",
            "mkdocs.yml": "site_name: Demo\ndocs_dir: handbook\n",
            "docusaurus.config.js": "module.exports = {};",
            "docs.json": '{"name":"Mintlify"}',
            "astro.config.mjs": "import starlight from '@astrojs/starlight';",
        }
        for path, content in files.items():
            self.write(path, content)
        self.track(*files)

        structures = self.scan()["structures"]

        self.assertEqual(structures["context"], ["CONTEXT.md"])
        self.assertEqual(structures["adr"], ["docs/adr/0001-use-sql.md"])
        self.assertEqual(
            [(site["kind"], site["path"]) for site in structures["docs_sites"]],
            [
                ("astro-starlight", "astro.config.mjs"),
                ("docusaurus", "docusaurus.config.js"),
                ("mintlify", "docs.json"),
                ("mkdocs", "mkdocs.yml"),
            ],
        )


class ManifestTests(SourceFixture):
    def test_manifest_bytes_are_stable_after_freshness_settles(self) -> None:
        self.write("docs/a.md", "alpha")
        self.track("docs/a.md")
        path = self.wiki / "_system" / "generated" / "sources" / "manifest.json"

        first = self.scan(generated_at="2026-09-05T01:00:00Z")
        pipeline.write_manifest(path, first, repo_root=self.root)
        second = self.scan(
            previous=first,
            generated_at="2026-09-05T02:00:00Z",
        )
        pipeline.write_manifest(path, second, repo_root=self.root)
        settled = path.read_bytes()
        third = self.scan(
            previous=second,
            generated_at="2026-09-05T03:00:00Z",
        )
        pipeline.write_manifest(path, third, repo_root=self.root)

        self.assertEqual(path.read_bytes(), settled)
        self.assertEqual(
            json.loads(settled)["generated_at"],
            "2026-09-05T02:00:00Z",
        )

    def test_reports_additions_modifications_and_deletions(self) -> None:
        self.write("docs/change.md", "before")
        self.write("docs/delete.md", "gone soon")
        self.track("docs/change.md", "docs/delete.md")
        previous = self.scan()
        self.write("docs/change.md", "after")
        (self.root / "docs/delete.md").unlink()
        subprocess.run(
            ["git", "rm", "--cached", "-q", "docs/delete.md"],
            cwd=self.root,
            check=True,
        )
        self.write("docs/add.md", "new")
        self.track("docs/add.md")

        current = self.scan(previous=previous)

        self.assertEqual(current["changes"]["added"], ["docs/add.md"])
        self.assertEqual(current["changes"]["modified"], ["docs/change.md"])
        self.assertEqual(current["changes"]["deleted"], ["docs/delete.md"])
        states = {entry["path"]: entry["freshness"] for entry in current["entries"]}
        self.assertEqual(states["docs/add.md"], "added")
        self.assertEqual(states["docs/change.md"], "modified")


class PointerTests(SourceFixture):
    def test_group_pointers_are_fixed_idempotent_and_do_not_write_semantic(self) -> None:
        self.write("CONTEXT.md", "# Context")
        self.write("docs/adr/0001-choice.md", "# Choice")
        self.write("mkdocs.yml", "site_name: Demo")
        self.write("data/input.csv", "a,b")
        self.track("CONTEXT.md", "docs/adr/0001-choice.md", "mkdocs.yml", "data/input.csv")
        manifest = self.scan()

        first = pipeline.write_raw_group_pointers(
            manifest,
            wiki_root=self.wiki,
            template_path=(
                self.wiki / "_system" / "templates" / "raw-group-pointer.md"
            ),
        )
        before = {path.name: path.read_bytes() for path in first["written"]}
        second = pipeline.write_raw_group_pointers(
            manifest,
            wiki_root=self.wiki,
            template_path=(
                self.wiki / "_system" / "templates" / "raw-group-pointer.md"
            ),
        )

        self.assertLessEqual(len(first["written"]), len(pipeline.POINTER_GROUPS))
        self.assertEqual(second["written"], first["written"])
        self.assertEqual(
            {path.name: path.read_bytes() for path in second["written"]},
            before,
        )
        self.assertFalse((self.wiki / "core").exists())

    def test_unmanaged_pointer_collision_is_preserved(self) -> None:
        self.write("CONTEXT.md", "# Context")
        self.track("CONTEXT.md")
        raw = self.wiki / "raw"
        raw.mkdir()
        collision = raw / "source-context.md"
        collision.write_text("human-owned", encoding="utf-8")

        result = pipeline.write_raw_group_pointers(
            self.scan(),
            wiki_root=self.wiki,
            template_path=(
                self.wiki / "_system" / "templates" / "raw-group-pointer.md"
            ),
        )

        self.assertEqual(collision.read_text(encoding="utf-8"), "human-owned")
        self.assertEqual(result["collisions"], [collision])


class RetrievalTests(SourceFixture):
    def test_raw_search_is_bounded_non_authoritative_and_skips_code(self) -> None:
        self.write(
            "docs/decision.md",
            "# Storage\n\nChoose durable quartz storage. " * 100,
        )
        self.write("src/quartz.py", "QUARTZ_SECRET_IMPLEMENTATION = True")
        self.track("docs/decision.md", "src/quartz.py")
        manifest = self.scan()

        results = pipeline.search_raw_sources(
            "quartz storage",
            manifest,
            repo_root=self.root,
            token_budget=50,
            per_result_tokens=35,
            max_results=3,
        )
        code_only = pipeline.search_raw_sources(
            "QUARTZ_SECRET_IMPLEMENTATION",
            manifest,
            repo_root=self.root,
        )

        self.assertEqual(len(results), 1)
        self.assertEqual(results[0]["provenance"]["authority"], "non-authoritative")
        self.assertEqual(results[0]["provenance"]["source_path"], "docs/decision.md")
        self.assertLessEqual(results[0]["tokens"], 35)
        self.assertLessEqual(sum(item["tokens"] for item in results), 50)
        self.assertEqual(code_only, [])

    def test_irrelevant_sources_are_not_returned(self) -> None:
        self.write("docs/cats.md", "feline whiskers")
        self.track("docs/cats.md")

        self.assertEqual(
            pipeline.search_raw_sources(
                "database quorum",
                self.scan(),
                repo_root=self.root,
            ),
            [],
        )


class ConversionCacheTests(SourceFixture):
    def test_cache_identity_reuse_and_source_invalidation(self) -> None:
        source = self.write("docs/spec.docx", "version one")
        self.track("docs/spec.docx")
        converter = MockConverter()
        first = self.scan(converter=converter)
        first_entry = first["entries"][0]
        second = self.scan(previous=first, converter=converter)
        self.assertEqual(len(converter.calls), 1)
        self.assertEqual(first_entry["conversion"]["state"], "converted")
        self.assertEqual(
            second["entries"][0]["conversion"]["state"],
            "cached",
        )
        first_key = first_entry["conversion"]["cache_key"]

        source.write_text("version two", encoding="utf-8")
        third = self.scan(previous=second, converter=converter)

        self.assertEqual(len(converter.calls), 2)
        self.assertNotEqual(
            third["entries"][0]["conversion"]["cache_key"],
            first_key,
        )

    def test_converter_config_changes_identity_and_failure_is_retryable(self) -> None:
        self.write("docs/spec.docx", "valid")
        self.write("docs/broken.docx", "invalid")
        self.track("docs/spec.docx", "docs/broken.docx")
        converter = MockConverter()
        one = self.scan(
            host={"sources": {"conversion": {"mode": "one"}}},
            converter=converter,
        )
        two = self.scan(
            host={"sources": {"conversion": {"mode": "two"}}},
            previous=one,
            converter=converter,
        )
        one_by_path = {entry["path"]: entry for entry in one["entries"]}
        two_by_path = {entry["path"]: entry for entry in two["entries"]}

        self.assertNotEqual(
            one_by_path["docs/spec.docx"]["conversion"]["cache_key"],
            two_by_path["docs/spec.docx"]["conversion"]["cache_key"],
        )
        failure = two_by_path["docs/broken.docx"]["conversion"]
        self.assertEqual(failure["state"], "failed")
        self.assertTrue(failure["retryable"])
        self.assertIn("malformed document", failure["diagnostic"])

    def test_derived_search_keeps_original_attribution(self) -> None:
        self.write("docs/spec.docx", "unique conversion phrase")
        self.track("docs/spec.docx")
        manifest = self.scan(converter=MockConverter())

        results = pipeline.search_raw_sources(
            "unique conversion",
            manifest,
            repo_root=self.root,
        )

        self.assertEqual(results[0]["provenance"]["source_path"], "docs/spec.docx")
        self.assertEqual(results[0]["provenance"]["content"], "derived")

    def test_compiled_claim_wins_and_conflict_is_idempotent(self) -> None:
        packed = [
            {
                "path": "core/existing-knowledge.md",
                "excerpt": "Do not use domain events.",
                "provenance": {"kind": "compiled"},
            },
            {
                "path": "docs/adr/0001-use-events.md",
                "excerpt": "Use domain events.",
                "provenance": {"kind": "raw"},
            },
        ]
        first = pipeline.detect_and_emit_conflicts(
            "domain events",
            packed,
            repo_root=self.root,
            wiki_root=self.wiki,
        )
        second = pipeline.detect_and_emit_conflicts(
            "domain events",
            packed,
            repo_root=self.root,
            wiki_root=self.wiki,
        )

        self.assertEqual(len(first), 1)
        self.assertEqual(first[0]["authoritative"], "core/existing-knowledge.md")
        self.assertEqual(first[0]["triage_path"], second[0]["triage_path"])
        inbox = list((self.wiki / "inbox").glob("*source-conflict*.md"))
        self.assertEqual(len(inbox), 1)


class ConversionPolicyTests(SourceFixture):
    def test_enabled_formats_honor_legacy_format_and_allowlist(self) -> None:
        self.assertEqual(
            pipeline.enabled_formats({"enabled": True, "format": "csv"}),
            ("csv",),
        )
        self.assertEqual(
            pipeline.enabled_formats({"enabled": True, "formats": ["CSV", "html"]}),
            ("csv", "html"),
        )
        self.assertEqual(
            pipeline.enabled_formats({"enabled": True}),
            pipeline.SAFE_LOCAL_FORMATS,
        )

    def test_disabled_formats_are_skipped_and_binaries_stay_inventoried(self) -> None:
        self.write("docs/keep.csv", "a,b\n1,2\n")
        self.write("docs/skip.pdf", b"%PDF-1.4\n")
        self.write("docs/blob.bin", b"\x00\x01\x02\x03")
        self.track("docs/keep.csv", "docs/skip.pdf", "docs/blob.bin")
        manifest = self.scan(
            host={
                "sources": {
                    "conversion": {"enabled": True, "formats": ["csv"]},
                }
            },
            converter=MockConverter(),
        )
        by_path = {entry["path"]: entry for entry in manifest["entries"]}
        self.assertEqual(by_path["docs/keep.csv"]["conversion"]["state"], "converted")
        self.assertEqual(by_path["docs/skip.pdf"]["conversion"]["state"], "skipped")
        self.assertEqual(by_path["docs/blob.bin"]["conversion"]["state"], "unsupported")

    def test_external_flags_block_conversion_without_allow_external(self) -> None:
        self.write("docs/keep.csv", "a,b\n1,2\n")
        self.track("docs/keep.csv")
        manifest = self.scan(
            host={
                "sources": {
                    "conversion": {
                        "enabled": True,
                        "formats": ["csv"],
                        "allow_urls": True,
                    }
                }
            },
            converter=MockConverter(),
        )
        conversion = manifest["entries"][0]["conversion"]
        self.assertEqual(conversion["state"], "blocked")
        self.assertFalse(conversion["retryable"])
        self.assertIn("allow_external", conversion["diagnostic"])

    def test_allow_external_is_required_before_external_flags(self) -> None:
        self.assertIsNotNone(
            pipeline.external_conversion_blocked({"allow_plugins": True})
        )
        self.assertIsNone(
            pipeline.external_conversion_blocked(
                {"allow_external": True, "allow_plugins": True}
            )
        )

    def test_commit_groups_copy_derived_markdown(self) -> None:
        self.write("docs/keep.csv", "a,b\n1,2\n")
        self.track("docs/keep.csv")
        manifest = self.scan(
            host={
                "sources": {
                    "conversion": {
                        "enabled": True,
                        "formats": ["csv"],
                        "commit_groups": ["data"],
                    }
                }
            },
            converter=MockConverter(),
        )
        committed = (
            self.root
            / "docs/wiki/_system/generated/sources/committed/docs/keep.csv.md"
        )
        self.assertTrue(committed.is_file())
        self.assertIn("Converted:", committed.read_text(encoding="utf-8"))
        self.assertEqual(manifest["entries"][0]["conversion"]["state"], "converted")

    def test_missing_format_extra_is_pending_and_retryable(self) -> None:
        self.write("docs/spec.pdf", b"%PDF-1.4\n")
        self.track("docs/spec.pdf")
        with mock.patch.object(
            pipeline,
            "format_extra_missing",
            return_value="MarkItDown extra `pdf` is not installed",
        ):
            manifest = self.scan(
                host={
                    "sources": {
                        "conversion": {"enabled": True, "formats": ["pdf"]},
                    }
                }
            )
        conversion = manifest["entries"][0]["conversion"]
        self.assertEqual(conversion["state"], "pending")
        self.assertTrue(conversion["retryable"])
        self.assertIn("pdf", conversion["diagnostic"])


def _minimal_pdf(text: str) -> bytes:
    stream = f"BT /F1 12 Tf 72 720 Td ({text}) Tj ET".encode("ascii")
    objects = [
        b"1 0 obj<< /Type /Catalog /Pages 2 0 R >>endobj\n",
        b"2 0 obj<< /Type /Pages /Kids [3 0 R] /Count 1 >>endobj\n",
        b"3 0 obj<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        b"/Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>endobj\n",
        b"4 0 obj<< /Length %d >>stream\n" % len(stream) + stream + b"\nendstream\nendobj\n",
        b"5 0 obj<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>endobj\n",
    ]
    body = b"".join(objects)
    offsets = []
    cursor = len(b"%PDF-1.4\n")
    for obj in objects:
        offsets.append(cursor)
        cursor += len(obj)
    xref = ["xref", "0 6", "0000000000 65535 f "]
    xref.extend(f"{offset:010d} 00000 n " for offset in offsets)
    startxref = len(b"%PDF-1.4\n") + len(body)
    return (
        b"%PDF-1.4\n"
        + body
        + ("\n".join(xref) + "\n").encode("ascii")
        + b"trailer<< /Size 6 /Root 1 0 R >>\nstartxref\n"
        + str(startxref).encode("ascii")
        + b"\n%%EOF\n"
    )


def _ooxml_zip(parts: dict[str, str]) -> bytes:
    import io
    import zipfile

    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w") as archive:
        for name, content in parts.items():
            archive.writestr(name, content)
    return buffer.getvalue()


def _minimal_docx(text: str) -> bytes:
    document = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
        f"<w:body><w:p><w:r><w:t>{text}</w:t></w:r></w:p></w:body></w:document>"
    )
    return _ooxml_zip(
        {
            "[Content_Types].xml": (
                '<?xml version="1.0" encoding="UTF-8"?>'
                '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
                '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
                '<Default Extension="xml" ContentType="application/xml"/>'
                '<Override PartName="/word/document.xml" '
                'ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>'
                "</Types>"
            ),
            "_rels/.rels": (
                '<?xml version="1.0" encoding="UTF-8"?>'
                '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
                '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>'
                "</Relationships>"
            ),
            "word/_rels/document.xml.rels": (
                '<?xml version="1.0" encoding="UTF-8"?>'
                '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"></Relationships>'
            ),
            "word/document.xml": document,
        }
    )


def _minimal_pptx(text: str) -> bytes:
    slide = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<p:sld xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" '
        'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
        "<p:cSld><p:spTree><p:nvGrpSpPr><p:cNvPr id=\"1\" name=\"\"/>"
        "<p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr><p:grpSpPr/>"
        f"<p:sp><p:nvSpPr><p:cNvPr id=\"2\" name=\"Title\"/><p:cNvSpPr/><p:nvPr/>"
        f"</p:nvSpPr><p:spPr/><p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:r><a:t>{text}</a:t></a:r></a:p>"
        "</p:txBody></p:sp></p:spTree></p:cSld></p:sld>"
    )
    return _ooxml_zip(
        {
            "[Content_Types].xml": (
                '<?xml version="1.0" encoding="UTF-8"?>'
                '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
                '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
                '<Default Extension="xml" ContentType="application/xml"/>'
                '<Override PartName="/ppt/presentation.xml" '
                'ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"/>'
                '<Override PartName="/ppt/slides/slide1.xml" '
                'ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/>'
                "</Types>"
            ),
            "_rels/.rels": (
                '<?xml version="1.0" encoding="UTF-8"?>'
                '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
                '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="ppt/presentation.xml"/>'
                "</Relationships>"
            ),
            "ppt/presentation.xml": (
                '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
                '<p:presentation xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
                '<p:sldIdLst><p:sldId id="256" r:id="rId1" '
                'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"/>'
                "</p:sldIdLst></p:presentation>"
            ),
            "ppt/_rels/presentation.xml.rels": (
                '<?xml version="1.0" encoding="UTF-8"?>'
                '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
                '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="slides/slide1.xml"/>'
                "</Relationships>"
            ),
            "ppt/slides/slide1.xml": slide,
        }
    )


def _minimal_xlsx(text: str) -> bytes:
    return _ooxml_zip(
        {
            "[Content_Types].xml": (
                '<?xml version="1.0" encoding="UTF-8"?>'
                '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
                '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
                '<Default Extension="xml" ContentType="application/xml"/>'
                '<Override PartName="/xl/workbook.xml" '
                'ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet.main+xml"/>'
                '<Override PartName="/xl/worksheets/sheet1.xml" '
                'ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.worksheet+xml"/>'
                '<Override PartName="/xl/sharedStrings.xml" '
                'ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.sharedStrings+xml"/>'
                "</Types>"
            ),
            "_rels/.rels": (
                '<?xml version="1.0" encoding="UTF-8"?>'
                '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
                '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="xl/workbook.xml"/>'
                "</Relationships>"
            ),
            "xl/workbook.xml": (
                '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
                '<workbook xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" '
                'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
                '<sheets><sheet name="Sheet1" sheetId="1" r:id="rId1"/></sheets></workbook>'
            ),
            "xl/_rels/workbook.xml.rels": (
                '<?xml version="1.0" encoding="UTF-8"?>'
                '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
                '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/worksheet" Target="worksheets/sheet1.xml"/>'
                '<Relationship Id="rId2" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/sharedStrings" Target="sharedStrings.xml"/>'
                "</Relationships>"
            ),
            "xl/worksheets/sheet1.xml": (
                '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
                '<worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main">'
                '<sheetData><row r="1"><c r="A1" t="s"><v>0</v></c></row></sheetData></worksheet>'
            ),
            "xl/sharedStrings.xml": (
                '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
                '<sst xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" count="1" uniqueCount="1">'
                f"<si><t>{text}</t></si></sst>"
            ),
        }
    )


def _minimal_epub(text: str) -> bytes:
    import io
    import zipfile

    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w") as archive:
        archive.writestr("mimetype", "application/epub+zip", compress_type=zipfile.ZIP_STORED)
        archive.writestr(
            "META-INF/container.xml",
            '<?xml version="1.0"?><container version="1.0" '
            'xmlns="urn:oasis:names:tc:opendocument:xmlns:container">'
            "<rootfiles><rootfile full-path=\"OEBPS/content.opf\" "
            'media-type="application/oebps-package+xml"/></rootfiles></container>',
        )
        archive.writestr(
            "OEBPS/content.opf",
            '<?xml version="1.0" encoding="UTF-8"?>'
            '<package xmlns="http://www.idpf.org/2007/opf" unique-identifier="bookid" version="2.0">'
            '<metadata xmlns:dc="http://purl.org/dc/elements/1.1/">'
            "<dc:identifier id=\"bookid\">repobrain-sample</dc:identifier>"
            "<dc:title>Sample</dc:title><dc:language>en</dc:language></metadata>"
            '<manifest><item id="n1" href="chapter.xhtml" media-type="application/xhtml+xml"/>'
            "</manifest><spine><itemref idref=\"n1\"/></spine></package>",
        )
        archive.writestr(
            "OEBPS/chapter.xhtml",
            '<?xml version="1.0" encoding="UTF-8"?>'
            '<html xmlns="http://www.w3.org/1999/xhtml"><head><title>Sample</title></head>'
            f"<body><p>{text}</p></body></html>",
        )
    return buffer.getvalue()


@unittest.skipUnless(
    pipeline.markitdown_info().get("compatible"),
    "markitdown==0.1.7 is required for live format tests",
)
class LiveMarkItDownFormatTests(SourceFixture):
    def test_live_local_formats_convert_and_keep_attribution(self) -> None:
        files: dict[str, str | bytes] = {
            "docs/sample.csv": "name,note\nlive-csv-token,ok\n",
            "docs/sample.html": "<html><body><p>live-html-token</p></body></html>",
            "docs/sample.epub": _minimal_epub("live-epub-token"),
            "docs/sample.pdf": _minimal_pdf("live-pdf-token"),
            "docs/sample.docx": _minimal_docx("live-docx-token"),
            "docs/sample.xlsx": _minimal_xlsx("live-xlsx-token"),
        }
        try:
            from pptx import Presentation

            pptx_path = self.root / "docs/sample.pptx"
            pptx_path.parent.mkdir(parents=True, exist_ok=True)
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])
            if slide.shapes.title:
                slide.shapes.title.text = "live-pptx-token"
            else:
                box = slide.shapes.add_textbox(0, 0, 1_000_000, 200_000)
                box.text_frame.text = "live-pptx-token"
            presentation.save(pptx_path)
            files["docs/sample.pptx"] = pptx_path.read_bytes()
        except ImportError:
            files["docs/sample.pptx"] = _minimal_pptx("live-pptx-token")
        for relative, content in files.items():
            self.write(relative, content)
        self.track(*files)
        converter = pipeline.live_converter()
        self.assertIsNotNone(converter)
        host = {
            "sources": {
                "conversion": {
                    "enabled": True,
                    "formats": list(pipeline.SAFE_LOCAL_FORMATS),
                }
            }
        }
        manifest = self.scan(host=host, converter=converter)
        by_path = {entry["path"]: entry for entry in manifest["entries"]}
        converted = 0
        for relative, token in (
            ("docs/sample.csv", "live-csv-token"),
            ("docs/sample.html", "live-html-token"),
            ("docs/sample.epub", "live-epub-token"),
            ("docs/sample.pdf", "live-pdf-token"),
            ("docs/sample.docx", "live-docx-token"),
            ("docs/sample.pptx", "live-pptx-token"),
            ("docs/sample.xlsx", "live-xlsx-token"),
        ):
            fmt = pipeline.format_for_suffix(Path(relative).suffix)
            extra = pipeline.format_extra_missing(fmt)
            conversion = by_path[relative]["conversion"]
            if extra:
                self.assertIn(conversion["state"], {"failed", "pending"}, relative)
                continue
            self.assertIn(
                conversion["state"],
                {"converted", "cached"},
                f"{relative}: {conversion.get('diagnostic')}",
            )
            converted += 1
            hits = pipeline.search_raw_sources(
                token,
                manifest,
                repo_root=self.root,
            )
            self.assertTrue(hits, relative)
            self.assertEqual(hits[0]["provenance"]["source_path"], relative)
            self.assertEqual(hits[0]["provenance"]["authority"], "non-authoritative")
        self.assertGreaterEqual(converted, 3)


if __name__ == "__main__":
    unittest.main()
